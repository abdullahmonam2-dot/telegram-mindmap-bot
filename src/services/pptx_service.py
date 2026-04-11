from pptx import Presentation
from pptx.util import Inches, Pt
from src.utils.logger import log
import os
import uuid
import re
from config import config

class PPTXService:
    async def create_presentation(self, slides_content: str):
        """Convert structured text to a PowerPoint file."""
        try:
            prs = Presentation()
            
            # Basic parsing of Gemini output
            # Expected format: Slide 1: Title \n Points...
            slides = re.split(r'الشريحة\s+\d+:|Slide\s+\d+:', slides_content)
            
            for slide_text in slides:
                if not slide_text.strip():
                    continue
                    
                lines = slide_text.strip().split('\n')
                title_text = lines[0].strip()
                body_text = "\n".join(lines[1:]).strip()
                
                slide_layout = prs.slide_layouts[1] # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = title_text
                
                body = slide.placeholders[1]
                body.text = body_text
                
            file_path = os.path.join(config.TEMP_DIR, f"slides_{uuid.uuid4()}.pptx")
            prs.save(file_path)
            log.info(f"PPTX generated: {file_path}")
            return file_path
        except Exception as e:
            log.error(f"PPTX error: {e}")
            return None

pptx_service = PPTXService()
